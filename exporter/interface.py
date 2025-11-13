
# ######################################### WORKING KODE


# from API2 import PresentationAPI
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches
# import torch
# # Importar las funciones del modelo
# from transformers import AutoModelForCausalLM, AutoTokenizer
# import io
# import gc
# #device = torch.device('cuda') 
# # Función para cargar el modelo
# @st.cache_resource
# def load_model(model_path):
#     gc.collect()
#     if torch.cuda.is_available():
#         torch.cuda.empty_cache()
#     model = AutoModelForCausalLM.from_pretrained(
#         model_path,
#         #device_map="auto",
#         trust_remote_code=True,
#         #torch_dtype="auto"
#         torch_dtype=torch.float16,#
#         #max_memory={0: "15GB"},#delete or edit according to the vram
#     )
#     tokenizer = AutoTokenizer.from_pretrained(
#         model_path,
#         trust_remote_code=True
#     )
#     return model, tokenizer

# # Function to generate a response
# def generate_response(instruction, model, tokenizer, max_length=512):
#     prompt = f"<|im_start|>human\n{instruction}<|im_end|>\n<|im_start|>assistant\n"
#     inputs = tokenizer(prompt, return_tensors="pt").to(model.device)
#     outputs = model.generate(
#         **inputs,
#         max_length=max_length,
#         pad_token_id=tokenizer.pad_token_id,
#         temperature=0.7,
#         do_sample=True,
#         top_p=0.95,
#     )
#     response = tokenizer.decode(outputs[0], skip_special_tokens=False)
#     return response

# ##############################################################################################################


# import streamlit as st
# from PIL import Image, ImageDraw
# import os
# from pdf2image import convert_from_path



# # API and required functions
# api = PresentationAPI('presentation.pptx')
# st.title("Chatbot Slide-Generator")

# # Load the model
# model_path = "/home/jovyan/2nd semester/PP Model/PerfectModel/Qwen2.5_Finetuned3.0"
# model, tokenizer = load_model(model_path)

# # Text entry area
# user_input = st.text_input("Type your prompt")




# def pptx_to_images(pptx_path, output_folder):
#     """
#     Convert each slide in a PowerPoint file to PNG images.

#     Args:
#         pptx_path (str): Path to the .pptx file.
#         output_folder (str): Folder where the images will be saved.

#     Returns:
#         List[str]: List of paths to the generated images.
#     """
#     presentation = Presentation(pptx_path)
#     slide_images = []

#     # Make sure the output folder exists
#     os.makedirs(output_folder, exist_ok=True)

#     for i, slide in enumerate(presentation.slides):
#         # Create a blank image for each slide
#         img = Image.new("RGB", (1280, 720), color="white")
#         draw = ImageDraw.Draw(img)
#         text = f"Slide {i + 1}"  
#         draw.text((50, 50), text, fill="black")  # Add example text
#         output_path = os.path.join(output_folder, f"slide_{i + 1}.png")
#         img.save(output_path)
#         slide_images.append(output_path)
    
#     return slide_images

# # Temporary folder to save the images
# output_folder = "slides_preview"





# # Button to generate the slide
# if st.button("Generate"):
#     if user_input:
#         with st.spinner("Generating response..."):
#             # Get the model response
#             response = generate_response(user_input, model, tokenizer)
            
#             # Create the slide with the answer
#             if "<|im_start|>assistant" in response:
#                 # Get content after "assistant"
#                 content = response.split("<|im_start|>assistant")[1]
#                 # Delete all other tags
#                 content = content.split("<|im_end|>")[0]
#                 content = content.split("<userStyle>")[0]
#                 content = content.split("<tool_call>")[0]
#                 content_before_tool_call = content.strip()
#             else:
#                 content_before_tool_call = response.strip()

#             # Show clean content
#             st.write(content_before_tool_call)

#             # If there is a <tool_call> execute it
#             if "<tool_call>" in response and "</tool_call>" in response:
#                     tool_call_code = response.split("<tool_call>")[1].split("</tool_call>")[0].strip()
                    
#                     # Show the code of the response
#                     #st.write("Generated code for tool_call:")
                    
#                     #st.code(tool_call_code)
                
#                     try:
#                         # Validate and run the generated code
#                         exec(tool_call_code)
                
#                         # Read the pptx generated file
#                         with open("presentation.pptx", "rb") as ppt_file:
#                             ppt_data = ppt_file.read()
                
#                         # Download the file
#                         st.success("Slide generated successfully!")
#                         st.download_button(
#                             label="Download Presentation",
#                             data=ppt_data,
#                             file_name="presentation.pptx",
#                             mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#                         )
#                     except SyntaxError as e:
#                         st.error(f"Syntax error in tool_call code: {e}")
#                     except Exception as e:
#                         st.error(f"Error while executing tool call: {e}")
                        
#                     else:
#                         st.error("Please enter a text.")
        
            
#             ####### preview if image slides (all white)
#             try:
#                 # Generar imágenes de las diapositivas
#                 slide_images = pptx_to_images("presentation.pptx", output_folder)

#                 # Mostrar las imágenes generadas en Streamlit
#                 for slide_image in slide_images:
#                     st.image(slide_image, caption=f"Slide {os.path.basename(slide_image)}", width=300)

#             except Exception as e:
#                 st.error(f"Error al generar imágenes: {e}")



# ######################################## WORKING KODE

from API2 import PresentationAPI
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import torch
from transformers import AutoModelForCausalLM, AutoTokenizer
import io
import gc
import streamlit as st
from PIL import Image, ImageDraw
import os
from pdf2image import convert_from_path


st.set_page_config(layout='wide') 
#device = torch.device('cuda') 
# Función para cargar el modelo
@st.cache_resource
def load_model(model_path):
    gc.collect()
    if torch.cuda.is_available():
        torch.cuda.empty_cache()
    model = AutoModelForCausalLM.from_pretrained(
        model_path,
        #device_map="auto",
        trust_remote_code=True,
        #torch_dtype="auto"
        torch_dtype=torch.float16,#
        #max_memory={0: "15GB"},#delete or edit according to the vram
    )
    tokenizer = AutoTokenizer.from_pretrained(
        model_path,
        trust_remote_code=True
    )
    return model, tokenizer

# Function to generate a response
def generate_response(instruction, model, tokenizer, max_length=2048):
    prompt = f"<|im_start|>human\n{instruction}<|im_end|>\n<|im_start|>assistant\n"
    inputs = tokenizer(prompt, return_tensors="pt").to(model.device)
    outputs = model.generate(
        **inputs,
        max_length=max_length,
        pad_token_id=tokenizer.pad_token_id,
        temperature=0.2,
        do_sample=True,
        top_p=0.95,
    )
    response = tokenizer.decode(outputs[0], skip_special_tokens=False)
    return response


def clean_up_generated_files():
    def on_rm_error(func, path, exc_info):
        """
        Function to delete .pptx file.
        """
        # allow permitions to delete
        os.chmod(path, stat.S_IWRITE)
        func(path)

    try:
        if os.path.exists("presentation.pptx"):
            os.remove("presentation.pptx")
            os.remove("presentation.pdf")
            os.system("cp template.pptx presentation.pptx")
            print("Deleted presentation.pptx")
        
        if os.path.exists("slide_images") and os.path.isdir("slide_images"):
            os.system("rm ./slides_preview/*.*")
            #shutil.rmtree("slide_images", onerror=on_rm_error)
            print("Deleted slide_images folder")
        
        st.success("Temporary files cleaned up successfully!")
    except Exception as e:
        st.error(f"Error while cleaning up files: {e}")

##############################################################################################################





# API and required functions
#api = PresentationAPI('presentation.pptx')

# Load the model
model_path = "/home/jovyan/PPTX/Qwen2.5_Fintuned3.0"
model, tokenizer = load_model(model_path)

# Temporary folder to save the images
output_folder = "slides_preview"

def pptx_to_images(pptx_path, output_folder):
    """
    Convert each slide in a PowerPoint file to PNG images.

    Args:
        pptx_path (str): Path to the .pptx file.
        output_folder (str): Folder where the images will be saved.

    Returns:
        List[str]: List of paths to the generated images.
    """
    presentation = Presentation(pptx_path)
    slide_images = []
    


    # Make sure the output folder exists
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, f"slide.jpg")
    os.system('libreoffice --headless --convert-to pdf ' + pptx_path)
    os.system('convert presentation.pdf ' + output_path)
    for i, slide in enumerate(presentation.slides):
        output_path = os.path.join(output_folder, f"slide-{i}.jpg")
        slide_images.append(output_path)   
    return slide_images

#    for i, slide in enumerate(presentation.slides):
        # Create a blank image for each slide
#        img = Image.new("RGB", (1280, 720), color="white")
#        draw = ImageDraw.Draw(img)
#        text = f"Slide {i + 1}"  
#        draw.text((50, 50), text, fill="black")  # Add example text
#        output_path = os.path.join(output_folder, f"slide_{i + 1}.png")
#        img.save(output_path)
#        slide_images.append(output_path)




################################### User Interface ###########################################################################

#create two colums for the layout desing
chat_col, preview_col=st.columns([0.7,0.3])

with preview_col:
    st.image("hof-logo.svg", width=200, channels="RGB", output_format="auto")
    
with chat_col:
   
    st.image("iisys-logo.svg", width=200, channels="RGB", output_format="auto")
    st.title("Presentation Slide Generator powered by LLM")
    # Text entry area
    #user_input = st.text_input("Type your prompt")
    user_input = st.text_area("Type your prompt", height=120)
    st.text("Tell the LLM what your presentation should be about, describe each slide or series of slides, explicitly mention how many bullet points you want, you can also add image slides. The current limit is six slides.")
    
    # Button to generate the slide
    if st.button("Generate"):
        if user_input:
            with st.spinner("Generating response..."):
                # Get the model response
                os.write(1, f"{user_input}\n".encode())
                response = generate_response(user_input, model, tokenizer)
                os.write(1, f"{response}\n".encode())
                # Create the slide with the answer
                if "<|im_start|>assistant" in response:
                    # Get content after "assistant"
                    content = response.split("<|im_start|>assistant")[1]
                    # Delete all other tags
                    content = content.split("<|im_end|>")[0]
                    content = content.split("<userStyle>")[0]
                    content = content.split("<tool_call>")[0]
                    content_before_tool_call = content.strip()
                else:
                    content_before_tool_call = response.strip()
    
                # Show clean content
                st.write(content_before_tool_call)
    
                # If there is a <tool_call> execute it
                if "<tool_call>" in response and "</tool_call>" in response:                        
                        tool_call_code = response.split("<tool_call>")[1].split("</tool_call>")[0].strip()
                        
                        # Show the code of the response
                        #st.write("Generated code for tool_call:")
                        
                        #st.code(tool_call_code)
                    
                        try:

                            #delete the old presentation and start with a fresh one from the defined template
                            #os.write(1, f"deleting previous files\n".encode())
                            #os.system("rm presentation.pptx")
                            #os.system("rm presentation.pdf")
                            #os.system("rm ./slides_preview/*.*")
                            #os.write(1, f"deleted previous files\n".encode())
                            #os.system("cp template.pptx presentation.pptx")

                            if 'ppt_data' not in st.session_state:
                                st.session_state['ppt_data'] = "value"
                            st.session_state.ppt_data = None 

                            api = PresentationAPI('presentation.pptx')

                            # Validate and run the generated code                            
                            os.write(1, f"Here comes the tool code\n".encode())
                            os.write(1, f"{tool_call_code}\n".encode())
                            exec(tool_call_code)
                            
                            # Read the pptx generated file
                            with open("presentation.pptx", "rb") as ppt_file:
                                ppt_data = ppt_file.read()
                    
                            # Download the file
                            st.success("Slide generated successfully!")
                            st.session_state.ppt_data = ppt_data  # Save presentation data in session state
                            st.session_state.generated = True  # Indicator that the presentation was generated
                            
                        except SyntaxError as e:
                            st.error(f"Syntax error in tool_call code: {e}")
      
 #-----------------------------visualization layout  

with preview_col:
    if st.session_state.get("generated", False):
        col1, col2 = st.columns(2)
        with col1:
            download_clicked = st.download_button(
                label="Download Presentation",
                data=st.session_state.ppt_data,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        with col2:
            if st.button("Clean Up Files"):
                clean_up_generated_files()
                st.session_state.generated = False
        
    if user_input and st.session_state.get("generated", False):
            ####### preview if image slides (all white)
            try:
                slide_images = pptx_to_images("presentation.pptx", output_folder)

                # Display images with a scroll container
                for slide_image in slide_images:
                    st.image(slide_image, caption=f"Slide {os.path.basename(slide_image)}", width=300)

            except Exception as e:
                st.error(f"Error generating images: {e}")




########################################



        

        
        
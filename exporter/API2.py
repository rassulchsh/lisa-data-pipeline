from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
#from imagegen import generate_image
from luisImageGen import generate_image
from pptx.shapes.picture import Picture


class PresentationAPI:
    def __init__(self, file_path='presentation.pptx'):
        """Initialize the PresentationAPI with the specified file path.
        
        Args:
            file_path (str): The path to the PowerPoint file. If it doesn't exist, a new presentation will be created.
        """
        self.file_path = file_path
        try:
            self.prs = Presentation(file_path)
        except Exception as e:
            self.prs = Presentation()
            print(f"Error loading presentation: {e}")
        
        self.current_slide = len(self.prs.slides) - 1 if len(self.prs.slides) > 0 else 0  

    def create_slide(self, layout='title', title='{}', content='{}', subtitle=None):
        """Create a new slide in the presentation and update the current slide index.
        
        Args:
            layout (str): The layout type for the slide.
            title (str): The title of the slide.
            content (str): The content of the slide.
            subtitle (str, optional): The subtitle of the slide.
        
        Returns:
            int: The index of the newly created slide.
        """
        layout_dict = {
            'title': 0, 'title and content': 1, 'section header': 2,
            'two content': 3, 'comparison': 4, 'title only': 5,
            'blank': 6, 'content with caption': 7, 'picture with caption': 8
        }
        slide_layout = self.prs.slide_layouts[layout_dict[layout.lower()]]
        slide = self.prs.slides.add_slide(slide_layout)

        if title != '{}':
            title_placeholder = slide.placeholders[0]
            title_placeholder.text = title

        if subtitle and layout.lower() == 'title':
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        if content != '{}' and layout.lower() != 'title':
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content

        self.save_presentation()
        self.current_slide = len(self.prs.slides) -1  # Update current slide index
        return self.current_slide

    def next_slide(self):
        """Move to the next slide in the presentation."""
        if self.current_slide < len(self.prs.slides)-1:
            self.current_slide += 1
        return self.current_slide

    def previous_slide(self):
        """Move to the previous slide in the presentation."""
        if self.current_slide > 0:
            self.current_slide -= 1
        return self.current_slide

    def get_current_slide(self):
        """Get the index of the current slide."""
        return self.current_slide

    def add_text(self, text, title_placeholder_num=0, content_placeholder_num=1):
        # """Add non-bullet text to the current slide.

        # Args:
        #     text (str): The text to add.
        #     title_placeholder_num (int): The index of the title placeholder (default is 0).
        #     content_placeholder_num (int): The index of the content placeholder (default is 1).
        # """
        # slide = self.prs.slides[self.current_slide]  # Use the current slide index

        # if title_placeholder_num < len(slide.placeholders):
        #     title_placeholder = slide.placeholders[title_placeholder_num]
        #     top_position = title_placeholder.top + title_placeholder.height
        # else:
        #     # Posición alternativa si el placeholder de título no está disponible
        #     top_position = Pt(100)  # Puedes ajustar esta posición según tus necesidades

        # # Intentar obtener el placeholder de contenido
        # if content_placeholder_num < len(slide.placeholders):
        #     content_placeholder = slide.placeholders[content_placeholder_num]
        #     left = content_placeholder.left
        #     width = content_placeholder.width
        #     height = content_placeholder.height
            
        # else:
        #     # Posición y tamaño alternativos si el placeholder de contenido no está disponible
        #     left = Pt(100)
        #     width = Pt(400)
        #     height = Pt(300)
        
        # textbox = slide.shapes.add_textbox(left, top_position, width, height)  # Ajusta la altura según tus necesidades
        # text_frame = textbox.text_frame
        # text_frame.word_wrap=True

        # p = text_frame.add_paragraph()
        # p.text = text
        # p.font.size = Pt(32)
        # p.alignment = PP_ALIGN.CENTER

        # self.save_presentation()

        slide = self.prs.slides[self.current_slide]
    
        try:
            # Obtener el placeholder de contenido y agregar el texto
            content_placeholder = slide.placeholders[content_placeholder_num]
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Limpiar contenido previo
            text_frame.word_wrap = True
            
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(32)
            p.alignment = PP_ALIGN.CENTER
            
            self.save_presentation()
            
        except Exception as e:
            print(f"Error adding text to placeholder: {str(e)}")






    
    # def add_bullet_points(self, slide_num, points, placeholder=1):
    #     """Add bullet points to a slide.

    #     Args:
    #         slide_num (int): The slide number to add bullet points to.
    #         points (list): A list of strings representing the bullet points.
    #         placeholder (int): The index of the placeholder for bullet points (default is 1).
    #     """
    #     if slide_num > len(self.prs.slides) or slide_num < 1:
    #         print(f"Error: slide {slide_num} does not exist.")
    #         return
    #     slide = self.prs.slides[slide_num]
    #     text_frame = slide.placeholders[placeholder].text_frame
    #     text_frame.clear()  # Clear previous content

    #     for point in points:
    #         p = text_frame.add_paragraph()  # Add a new paragraph
    #         p.text = point
    #         p.level = 0  # Set bullet level
    #     self.save_presentation()

    def add_bullet_points(self, slide_num, points, placeholder=1):
        """Add bullet points to a slide.

        Args:
            slide_num (int): The slide number to add bullet points to.
            points (list): A list of strings representing the bullet points.
            placeholder (int): The index of the placeholder for bullet points (default is 1).
        """
        if slide_num > len(self.prs.slides) or slide_num < 1:
            print(f"Error: slide {slide_num} does not exist.")
            return
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        text_frame.clear()  # Clear previous content

        for point_str in points:
            
            if '(' in point_str and ')' in point_str and ',' in point_str:
            
                last_comma = point_str.rindex(',')
                text = point_str[1:last_comma].strip()  # El [1:] remueve el primer paréntesis
                #print(text)
                level = int(point_str[last_comma + 1:-1].strip())
                
            else:
                text = point_str
                level = 0
            p = text_frame.add_paragraph()
            p.text = text
            p.level = level # Set bullet level
        self.save_presentation()





    def edit_text(self, slide_num, placeholder, new_text, paragraph_index=0):#creates a * to the left
        """Edit existing text in a slide's placeholder.

        Args:
            slide_num (int): The slide number to edit.
            placeholder (int): The index of the placeholder to edit.
            new_text (str): The new text to replace the existing text.
            paragraph_index (int): The index of the paragraph to edit (default is 0).
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        if paragraph_index < len(text_frame.paragraphs):
            text_frame.paragraphs[paragraph_index].text = new_text
            self.save_presentation()
        else:
            print(f"Paragraph index {paragraph_index} does not exist in slide {slide_num}, placeholder {placeholder}")

    def delete_text(self, slide_num, placeholder, paragraph_index=None):
        """Delete text from a slide's placeholder.

        Args:
            slide_num (int): The slide number to delete text from.
            placeholder (int): The index of the placeholder to delete text from.
            paragraph_index (int, optional): The index of the paragraph to delete. If None, all text will be deleted.
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        if paragraph_index is None:
            text_frame.clear()
        elif paragraph_index < len(text_frame.paragraphs):
            text_frame._element.remove(text_frame.paragraphs[paragraph_index]._element)
        else:
            print(f"Paragraph index {paragraph_index} does not exist in slide {slide_num}, placeholder {placeholder}")
        self.save_presentation()

    def change_font(self, slide_num, placeholder_num, font_name, font_size):
        """Change the font of text in a slide's placeholder.

        Args:
            slide_num (int): The slide number to change the font.
            placeholder_num (int): The index of the placeholder to change the font.
            font_name (str): The name of the font to apply.
            font_size (int): The size of the font to apply.
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder_num].text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
        self.save_presentation()

    def change_background_color(self, slide_num, color):
        """Change the background color of a slide.

        Args:
            slide_num (int): The slide number to change the background color.
            color (tuple): A tuple of RGB values (R, G, B) for the background color.
        """
        slide = self.prs.slides[slide_num]
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)
        self.save_presentation()

    def save_presentation(self):
        """Save the presentation to the specified file path."""
        self.prs.save(self.file_path)
        print(f"Presentation saved as '{self.file_path}'")

    def get_slide_count(self):
        """Get the total number of slides in the presentation.

        Returns:
            int: The total number of slides.
        """
        return len(self.prs.slides)

    def add_image(self, slide_num, prompt, image_path='generated_image.png'):
        """Add an image to a slide.

        Args:
            slide_num (int): The slide number to add the image to.
            prompt (str): The prompt for generating the image.
            image_path (str): The path where the generated image is saved.
        """
        # Assuming generate_image saves the image to image_path
        image_generated = generate_image(prompt)

        # Check if the image was generated successfully
        if image_generated:
            # Add the image to the slide
            slide = self.prs.slides[slide_num]

            # Position and size adjustments for the image
            left = Inches(5)
            top = Inches(2.5)
            width = Inches(4)
            height = Inches(3)

            # Insert the image into the slide
            slide.shapes.add_picture(image_path, left, top, width, height)

            # Save the presentation
            self.save_presentation()

    def delete_images(self, slide_num):
        """Delete all images from a slide.

        Args:
            slide_num (int): The slide number to delete images from.
        """
        slide = self.prs.slides[slide_num]
        # Create a list of shapes that are pictures
        images_to_delete = [shape for shape in slide.shapes if isinstance(shape, Picture)]

        # Remove the images from the slide
        for image in images_to_delete:
            sp = image._element
            sp.getparent().remove(sp)

        self.save_presentation()

    def add_picture_with_caption(self, slide_num, prompt, title, text, image_path='generated_image.png'):
        """
        Add a generated picture with a caption to a slide based on a prompt.

        Args:
            slide_num (int): The slide number where the picture and caption will be added.
            prompt (str): The prompt to generate the image.
            caption_text (str): The caption text to add below the image.
        """
        if slide_num >= len(self.prs.slides) or slide_num < 0:
            print(f"Error: slide {slide_num} does not exist.")
            return
        
        image_generated=generate_image(prompt)

        if image_generated:
        # Add the image to the slide
            slide = self.prs.slides[slide_num]

            image_placeholder = slide.shapes[1]
            image_placeholder.fill.solid()
            left = image_placeholder.left
            top = image_placeholder.top
            width = image_placeholder.width
            height = image_placeholder.height
            slide.shapes.add_picture(image_path, left, top, width, height)

            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            caption_placeholder = slide.placeholders[2]  
            caption_placeholder.text = text

        self.save_presentation()


    def delete_slide(self, slide_num):
        """Delete a slide from the presentation.
        
        Args:
            slide_num (int): The index of the slide to delete.
        """
        xml_slides = self.prs.slides._sldIdLst  # Lista de IDs de diapositivas
        slides = list(xml_slides)  # Convertimos a lista para manipular

        if 0 <= slide_num < len(slides):
            xml_slides.remove(slides[slide_num])  # Eliminamos la diapositiva
            self.prs.save(self.file_path)
        
                
        self.save_presentation()
        print(f"Slide {slide_num} deleted successfully.")



    def save_presentation(self):
        """Save the presentation to the specified file path."""
        self.prs.save(self.file_path)
        print(f"Presentation saved as '{self.file_path}'")

    # Add methods for other functionalities (add_bullet_points, edit_text, etc.)



# Example usage
#api = PresentationAPI('presentation.pptx')

# points = [
#     ('Fastest mode of transportation for long distances', 0),
#     ('Reduces travel time significantly', 1),
#     ('Connects remote and inaccessible areas', 0),
#     ('Enables access to islands and mountainous regions', 1),
#     ('Efficient for international travel', 0),
#     ('Supports global mobility for businesses and individuals', 1),
#     ('High safety standards compared to other modes of transport', 0),
#     ('Advanced technology and strict regulations ensure safety', 1),
#     ('Supports global trade and tourism', 0),
#     ('Facilitates the export and import of goods', 1),
#     ('Boosts local economies through tourism', 2)
# ]



# # Llamar a la función de agregar bullet points
# api.add_bullet_points(api.get_current_slide(), points=points, placeholder=1)
